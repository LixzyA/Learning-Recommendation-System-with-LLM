import csv
from pptx import Presentation
import docx2txt
from PyPDF2 import PdfReader
import os
import cv2
from langdetect import detect
import easyocr

def get_file_type(filename:str): 
    if filename.endswith('.pptx'):
        return "pptx"
    elif filename.endswith('.docx'):
        return "docx"
    elif filename.endswith(".txt"):
        return "txt"
    elif filename.endswith(".md"):
        return 'md'
    elif filename.endswith(".jpg"):
        return "jpg"
    elif filename.endswith(".jpeg"):
        return 'jpeg'
    elif filename.endswith('png'):
        return 'png'
    elif filename.endswith(".pdf"):
        return 'pdf'
    return 'unknown'

def get_lang(content:str):
    return detect(content)

if __name__ == "__main__":
    reader = easyocr.Reader(['en', 'ch_sim'], gpu=False)

    # get list of file in dir
    path = "dataset/"
    filename = 'supervisor-dataset-new.csv'

    files = os.listdir(path)

    # Open CSV file for writing
    with open(filename, "w", newline='', encoding="utf-8") as dataset:
        writer = csv.writer(dataset)
        writer.writerow(["name", "lang", "file_type" ,"content"])

        for file in files:
            text = ""  # Reset text for each file
            file = file.encode('utf-8').decode('utf-8')
            file_type = get_file_type(file)

            if file_type == 'pptx':
                try:
                    presentation = Presentation(path + file)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text
                    text = text.replace('\n', ' ')
                    

                except Exception as e:
                    print(f"Error processing PPTX {file}: {e}")

            elif file_type == 'docx':
                try:
                    text = docx2txt.process(path+file)
                    text = text.replace('\n', ' ')
                except Exception as e:
                    print(f"Error processing DOCX {file}: {e}")

            elif file_type == 'pdf':
                try:
                    pdf = PdfReader(path + file)
                    for page in pdf.pages:
                        text += page.extract_text() or ''
                    text = text.replace('\n', ' ')
                except Exception as e:
                    print(f"Error processing PDF {file}: {e}")

            elif file_type == 'txt' or file_type == 'md':
                try:
                    with open(path + file, encoding='utf-8') as txt:
                        text = txt.read().replace('\n', ' ')
                except UnicodeDecodeError:
                    print(f"Could not decode {file} with UTF-8.")

            elif file_type in ['jpg', 'jpeg', 'png']:
                try:
                    img = cv2.imread(path+file)
                    text_detections = reader.readtext(img)
                    threshold = 0.25

                    for texts in text_detections:
                        if texts[2] > threshold:
                            text+= texts[1] + " "

                    text = text.replace('\n', ' ')
                    
                except Exception as e:
                    print(f"Error processing img {file}: {e}")

            else:
                print(f"Unsupported file type: {file}")

            lang = get_lang(text)
            writer.writerow([file, lang, file_type,text])
